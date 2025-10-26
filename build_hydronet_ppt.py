# -*- coding: utf-8 -*-
"""Generate the HydroNet presentation deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

COLOR_BG_DARK = RGBColor(0, 45, 114)
COLOR_BG_BLUE = RGBColor(0, 88, 198)
COLOR_CYAN = RGBColor(0, 173, 239)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_SILVER = RGBColor(220, 230, 245)

TITLE_FONT = "微软雅黑"
BODY_FONT = "微软雅黑"


def add_bg_gradient(slide, top_color=COLOR_BG_DARK, bottom_color=COLOR_BG_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    fill = shape.fill
    fill.gradient()
    gradient = fill.gradient()
    gradient.stops[0].position = 0.0
    gradient.stops[0].color.rgb = top_color
    gradient.stops[1].position = 1.0
    gradient.stops[1].color.rgb = bottom_color
    shape.line.fill.background()


def add_title(
    slide,
    text,
    font_size=40,
    color=COLOR_WHITE,
    left=Inches(0.7),
    top=Inches(0.5),
    width=Inches(12.5),
    height=Inches(1.2),
    align=PP_ALIGN.LEFT,
    bold=True,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = TITLE_FONT
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    paragraph.alignment = align
    return textbox


def add_bullets(
    slide,
    items,
    left=Inches(0.9),
    top=Inches(1.9),
    width=Inches(12.0),
    height=Inches(4.8),
    font_size=24,
    color=COLOR_WHITE,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    for index, line in enumerate(items):
        paragraph = text_frame.add_paragraph() if index > 0 else text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = line
        font = run.font
        font.name = BODY_FONT
        font.size = Pt(font_size)
        font.color.rgb = color
        paragraph.level = 0
    return textbox


def add_caption(
    slide,
    text,
    left=Inches(0.9),
    top=Inches(6.9),
    width=Inches(12.0),
    height=Inches(0.7),
    font_size=14,
    color=COLOR_SILVER,
    align=PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = BODY_FONT
    font.size = Pt(font_size)
    font.color.rgb = color
    paragraph.alignment = align
    return textbox


def set_notes(slide, lines):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    for index, line in enumerate(lines):
        paragraph = notes_frame.add_paragraph() if index > 0 else notes_frame.paragraphs[0]
        paragraph.text = line


def add_chapter_cover(prs, chapter_number, chapter_title, tagline="南水北调集团 · HydroNet 发布会 2025"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_gradient(slide)
    add_title(
        slide,
        f"第{chapter_number}章  {chapter_title}",
        font_size=46,
        align=PP_ALIGN.CENTER,
        left=Inches(0.5),
        top=Inches(2.2),
        width=Inches(13.0),
    )
    add_caption(
        slide,
        tagline,
        align=PP_ALIGN.CENTER,
        left=Inches(0.5),
        top=Inches(3.5),
        width=Inches(13.0),
    )
    set_notes(
        slide,
        [
            f"章节封面：第{chapter_number}章 {chapter_title}",
            "动画建议：标题淡入（0.7s），副标题延迟淡入（0.4s）。",
        ],
    )
    return slide


def add_std_slide(prs, title, bullets, notes=None, caption="南水北调集团 · HydroNet 发布会 2025"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_gradient(slide)
    add_title(slide, title)
    add_bullets(slide, bullets)
    add_caption(slide, caption)
    set_notes(
        slide,
        notes
        or [
            "讲解要点：逐条说明。",
            "动画建议：项目符号按段落依次出现（0.3s阶梯）。",
        ],
    )
    return slide


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_width, slide_height = prs.slide_width, prs.slide_height

cover = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(cover)
add_title(
    cover,
    "水网大模型与自主运行水网智能体体系（HydroNet）",
    font_size=44,
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(1.9),
    width=Inches(12.8),
)
add_title(
    cover,
    "HydroNet — The Autonomous Hydro-System Model of China",
    font_size=22,
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(2.9),
    width=Inches(12.8),
    bold=False,
)
add_caption(
    cover,
    "让每一滴水，都被智能而高效地调控",
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(3.6),
    width=Inches(12.8),
)
add_caption(
    cover,
    "河北工程大学 · 智慧水网创新团队 | 主讲人：雷晓辉 | 2025",
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(6.9),
    width=Inches(12.8),
)
set_notes(
    cover,
    [
        "封面讲解：一句话定义HydroNet，突出‘认知+具身+HydroOS’。",
        "动画建议：主标题放大渐显，口号从下方上移淡入。",
    ],
)

add_std_slide(
    prs,
    "目录",
    [
        "一、理论与理念",
        "二、HydroNet总体架构",
        "三、技术体系与创新",
        "四、HydroAgent管理建模平台",
        "五、应用案例",
        "六、人才体系与未来蓝图",
        "团队与致谢",
    ],
    notes=["目录作为导航，快速预览各章重点。", "动画建议：整页淡入。"],
)

add_chapter_cover(prs, 1, "理论与理念")
slides_1 = [
    (
        "从传统水利到无人驾驶水网",
        [
            "阶段演进：传统水利 → 数字孪生 → 无人驾驶水网",
            "目标：从开环可视化走向闭环自治控制",
            "价值：效率、安全、生态、经济四维提升",
        ],
        ["对比‘会说话AI’与‘会调水AI’。"],
    ),
    (
        "水系统控制论的闭环",
        [
            "闭环：感知 → 预测 → 调度 → 控制 → 验证",
            "关键：模型、反馈、稳定性",
            "方法：前馈+反馈、鲁棒性与可解释性",
        ],
        None,
    ),
    (
        "快思慢想与云边协同",
        [
            "云端大脑（慢思）：MPC、优化、策略规划",
            "边端小脑（快思）：PID、反馈、毫秒级响应",
            "二者联动：目标分解、在环验证、闭环收敛",
        ],
        None,
    ),
    (
        "具身智能体：大脑与小脑",
        [
            "大脑：中央调度智能体（HydroBrain）",
            "小脑：本地控制智能体（HydroCerebellum）",
            "身体：数字孪生体（HydroTwin）",
        ],
        None,
    ),
    (
        "算法→算子→模型→智能体",
        [
            "算法沉淀为可复用算子",
            "算子装配形成可解释模型",
            "模型容器化为智能体服务",
        ],
        None,
    ),
    (
        "开环 vs 闭环",
        [
            "开环：方案→下发→缺乏自适应反馈",
            "闭环：反馈校正→在线辨识→持续优化",
            "指标：稳定性、超调量、能耗、可靠性",
        ],
        None,
    ),
    (
        "多学科融合",
        [
            "水文学、水力学、运筹学、控制论、人工智能",
            "数据同化：传感器+遥测+气象融合",
            "评估：鲁棒性、可用性、可维护性",
        ],
        None,
    ),
    (
        "无人驾驶水网定位",
        [
            "目标：自主感知、自主决策、自主控制",
            "路径：X-in-Loop在环测试体系保障",
            "成果：从‘智慧水利’迈向‘自主水网’",
        ],
        None,
    ),
    (
        "理念小结",
        [
            "闭环是核心、具身是关键、OS是保障",
            "语言即调度、模型即控制、仿真即验证",
            "HydroNet让AI从理解世界到改变世界",
        ],
        None,
    ),
]
for title, bullets, notes in slides_1:
    add_std_slide(prs, title, bullets, notes)

add_chapter_cover(prs, 2, "HydroNet总体架构")
slides_2 = [
    (
        "五层结构总览",
        [
            "认知：HydroGPT（理解、规划、任务生成）",
            "调度：HydroBrain（预测、优化、预演）",
            "控制：HydroCerebellum（执行、反馈）",
            "孪生：HydroTwin（感知、仿真、验证）",
            "OS：HydroOS（通信、编排、安全）",
        ],
        None,
    ),
    (
        "HydroGPT：认知智能体",
        [
            "自然语言→任务契约→调用具身服务",
            "知识图谱+RAG+Prompt 链",
            "不直接控设备，走预演安全闸",
        ],
        None,
    ),
    (
        "HydroBrain：中央调度智能体",
        [
            "预测（入流/需水/风险）+ MPC优化",
            "多场景预演与目标分解",
            "分布式孪生反馈智能体：在线辨识、降阶",
        ],
        None,
    ),
    (
        "HydroCerebellum：本地控制智能体",
        [
            "泵站经济运行、PID/自适应控制",
            "边缘一体机、毫秒级快思考",
            "目标跟踪与稳定性保障",
        ],
        None,
    ),
    (
        "HydroTwin：数字孪生体",
        [
            "施垚（产汇流/参数分区），王孝群（1D水动力）",
            "张雷/陈凯歌（传感器/执行器仿真）",
            "SIL/HIL验证与虚实映射",
        ],
        None,
    ),
    (
        "HydroOS：智能操作系统",
        [
            "分布式异构多智能体协同协议",
            "任务编排、模型注册、安全治理",
            "云-边-端一体化运行",
        ],
        None,
    ),
    (
        "层间数据流动",
        [
            "认知→调度：任务契约与约束",
            "调度→控制：目标轨迹与容差",
            "孪生→各层：状态与评估",
        ],
        None,
    ),
]
for title, bullets, notes in slides_2:
    add_std_slide(prs, title, bullets, notes)

add_chapter_cover(prs, 3, "技术体系与创新")
slides_3 = [
    (
        "分层分布式具身智能体体系",
        [
            "云：全局优化与策略规划",
            "边：本地快思反馈",
            "端：被控对象的感知与执行",
        ],
        None,
    ),
    (
        "云边协同：快思慢想",
        [
            "中心MPC（慢思） + 现场PID（快思）",
            "扰动抑制、目标分解、容错切换",
            "异常→限域→降级→恢复",
        ],
        None,
    ),
    (
        "X-in-Loop在环测试",
        [
            "MIL 模型在环",
            "SIL 软件在环",
            "HIL 硬件在环 / HITL 人在环 / SITL 系统在环",
        ],
        ["强调‘能用、好用、可靠’三层验证。"],
    ),
    (
        "通用仿真模型框架",
        [
            "水文/水力/水质模型库装配",
            "多拓扑（河道/渠道/管网/梯级）",
            "高保真+降阶并存",
        ],
        None,
    ),
    (
        "通用调度模型框架",
        [
            "统一建模 → 通用求解 → 知识决策",
            "多目标（供水/防洪/能耗/生态）",
            "鲁棒优化与方案评估",
        ],
        None,
    ),
    (
        "在线辨识与反馈",
        [
            "参数/状态在线估计",
            "模型漂移监测与自动校正",
            "AIOps：性能对比与回滚",
        ],
        None,
    ),
    (
        "数据同化与预测校正",
        [
            "传感器+遥测+气象融合",
            "短中期耦合预报",
            "偏差更正与不确定性量化",
        ],
        None,
    ),
    (
        "稳定性与鲁棒性",
        [
            "超调量/稳态误差/KPI，对比基线",
            "极端扰动情景压力测试",
            "冗余与主备切换",
        ],
        None,
    ),
    (
        "降阶建模与LQR/LQG",
        [
            "ID/IDZ/Muskingum-Cunge等",
            "线性化可控性/可观性分析",
            "控制律与代价函数设计",
        ],
        None,
    ),
]
for title, bullets, notes in slides_3:
    add_std_slide(prs, title, bullets, notes)

add_chapter_cover(prs, 4, "HydroAgent管理建模平台")
slides_4 = [
    (
        "平台结构",
        [
            "核心：HydroOS",
            "工作台A：认知智能体Studio",
            "工作台B：具身智能体Studio",
        ],
        None,
    ),
    (
        "LLM辅助建模",
        [
            "自然语言输入→自动生成流程图",
            "任务契约与参数补全",
            "一键预演与报告生成",
        ],
        None,
    ),
    (
        "拖拽式建模",
        [
            "组件化拼装：水文/水力/预测/控制",
            "参数面板/版本管理/回放对比",
            "模板沉淀与复用",
        ],
        None,
    ),
    (
        "注册中心与模型仓库",
        [
            "智能体注册、健康监测、依赖管理",
            "模型/算子版本追踪与可复现",
            "数据血缘与合规模块",
        ],
        None,
    ),
    (
        "分布式计算编排",
        [
            "批/流/实时融合",
            "云-边-端协同",
            "容错、断点续跑、优先级",
        ],
        None,
    ),
    (
        "场景库与预演",
        [
            "雨量站动态图、极端场景、故障注入",
            "多目标权衡与鲁棒性评估",
            "人机共决与安全闸",
        ],
        None,
    ),
    (
        "安全治理机制",
        [
            "多租户、最小权限、审计追踪",
            "算子白名单与双人复核",
            "‘语言即调度’的二级闸",
        ],
        None,
    ),
]
for title, bullets, notes in slides_4:
    add_std_slide(prs, title, bullets, notes)

add_chapter_cover(prs, 5, "应用案例")
slides_5 = [
    (
        "胶东调水工程（右脑快思）",
        [
            "多模态需求→知识图谱→最优化问题",
            "自动生成情景方案与多目标权衡",
            "方案工具调用→仿真评估→故事化响应",
        ],
        None,
    ),
    (
        "大渡河沙坪二级水电站（小脑控制）",
        [
            "泵站经济运行+PID实时控制",
            "越线风险/频繁扰动情景抑制",
            "KPI：能效/响应/稳定性提升",
        ],
        None,
    ),
    (
        "水库群联合调度",
        [
            "跨流域耦合目标（供水/防洪/发电/生态）",
            "MPC+鲁棒优化",
            "策略分解与跟踪",
        ],
        None,
    ),
    (
        "城市供排水数字孪生",
        [
            "一体化厂网河湖系统",
            "污染/洪峰/设备故障情景",
            "在环验证与应急联动",
        ],
        None,
    ),
    (
        "多源数据融合",
        [
            "气象/雷达/遥测/传感器同化",
            "短临预报联动控制",
            "自适应阈值与告警",
        ],
        None,
    ),
    (
        "成效量化",
        [
            "效率↑ 能耗↓ 可靠性↑",
            "样板工程可复制推广",
            "与监管指标对齐",
        ],
        None,
    ),
    (
        "案例小结",
        [
            "认知与具身闭环的工程化落地",
            "平台化与模板化复用",
            "体系化能力构建",
        ],
        None,
    ),
]
for title, bullets, notes in slides_5:
    add_std_slide(prs, title, bullets, notes)

add_chapter_cover(prs, 6, "人才体系与未来蓝图")
slides_6 = [
    (
        "课程与教材体系",
        [
            "《水系统感知/物联/仿真/辨识/控制/测试》",
            "从仿真导向到控制导向的转型",
            "实验平台与竞赛体系",
        ],
        None,
    ),
    (
        "复合型人才培养",
        [
            "水利+控制+AI+机电+通信",
            "工程导向与科研导向并重",
            "产业导师与实战项目",
        ],
        None,
    ),
    (
        "AI革命与工业革命",
        [
            "工业革命：有没有 → AI革命：好不好",
            "水网：强身健体 → 无人驾驶",
            "面向自治的工程范式变革",
        ],
        None,
    ),
    (
        "全球智慧水网生态",
        [
            "多流域多目标协同",
            "跨行业（水-能-城-航运）联动",
            "开放标准与国际合作",
        ],
        None,
    ),
    (
        "团队与致谢",
        [
            "河北工程大学 · 智慧水网创新团队",
            "研究方向：HydroNet/多智能体/在环测试",
            "合作单位与生态伙伴",
        ],
        None,
    ),
]
for title, bullets, notes in slides_6:
    add_std_slide(prs, title, bullets, notes)

end_slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(end_slide)
add_title(
    end_slide,
    "谢谢！",
    font_size=44,
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(2.2),
    width=Inches(12.8),
)
add_caption(
    end_slide,
    "让每一滴水，都被智能而高效地调控",
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(3.4),
    width=Inches(12.8),
)
add_caption(
    end_slide,
    "南水北调集团 · HydroNet 发布会 2025",
    align=PP_ALIGN.CENTER,
    left=Inches(0.4),
    top=Inches(6.9),
    width=Inches(12.8),
)
set_notes(
    end_slide,
    ["封底：简洁致谢。", "动画建议：主标题渐显，口号光晕淡入。"],
)

OUTPUT_NAME = "水网大模型与自主运行水网智能体体系（HydroNet）.pptx"
prs.save(OUTPUT_NAME)
print(f"已生成：{OUTPUT_NAME}")
